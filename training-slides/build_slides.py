# -*- coding: utf-8 -*-
"""Tao slide training Team Sales - tone mau & font thuong hieu Kingsmen (keokingsmen.com)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- Bang mau thuong hieu Kingsmen ----
INK       = RGBColor(0x0b, 0x35, 0x43)  # teal dam (chinh)
INK_SOFT  = RGBColor(0x11, 0x46, 0x54)
INK_MUTED = RGBColor(0x5c, 0x74, 0x80)
BRAND     = RGBColor(0x0a, 0x92, 0xb4)  # cyan thuong hieu
BRAND_DK  = RGBColor(0x0a, 0x6a, 0x80)
BRAND_LT  = RGBColor(0x7f, 0xd4, 0xe4)
BRAND_BG  = RGBColor(0xea, 0xfa, 0xfd)
LINE      = RGBColor(0xe2, 0xe8, 0xea)
PAGE_BG   = RGBColor(0xf6, 0xf7, 0xf9)
WHITE     = RGBColor(0xff, 0xff, 0xff)

F_HEAD = "Montserrat"
F_BODY = "Maven Pro"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color


def rect(s, x, y, w, h, color, shape=MSO_SHAPE.RECTANGLE, line_col=None, line_w=None):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line_col is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_col
        sp.line.width = line_w or Pt(1)
    sp.shadow.inherit = False
    return sp


def grad(s, x, y, w, h, c1, c2, angle=45):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sp.line.fill.background()
    sp.shadow.inherit = False
    sp.fill.gradient()
    stops = sp.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[1].position = 1.0
    stops[1].color.rgb = c2
    try:
        sp.fill.gradient_angle = angle
    except Exception:
        pass
    return sp


def txt(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold, font, spacing_override)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (text, size, color, bold, font, *rest) in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
            if rest and rest[0] is not None:
                # letter spacing (in points*100) via xml
                rPr = r._r.get_or_add_rPr()
                rPr.set('spc', str(int(rest[0] * 100)))
    return tb


def chip(s, x, y, w, text, fill, tcol, size=12):
    c = rect(s, x, y, w, Inches(0.42), fill, MSO_SHAPE.ROUNDED_RECTANGLE)
    try:
        c.adjustments[0] = 0.5
    except Exception:
        pass
    tf = c.text_frame
    tf.word_wrap = True
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = F_HEAD
    r.font.color.rgb = tcol
    return c


def page_footer(s, n, section):
    # thanh chan trang
    txt(s, Inches(0.75), Inches(6.95), Inches(6), Inches(0.4),
        [[("KINGSMEN", 10, BRAND_DK, True, F_HEAD, 2.5),
          ("  •  Seeding Marketing → Sales", 10, INK_MUTED, False, F_BODY)]],
        anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(11.0), Inches(6.95), Inches(1.6), Inches(0.4),
        [[(f"{section}   ", 10, INK_MUTED, False, F_BODY),
          (f"0{n}", 12, BRAND_DK, True, F_HEAD)]],
        align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(s, kicker, title, section, n):
    rect(s, 0, 0, Inches(0.28), EMU_H, BRAND)                 # thanh doc trai
    rect(s, Inches(0.75), Inches(0.62), Inches(0.55), Inches(0.11), BRAND)
    txt(s, Inches(0.75), Inches(0.75), Inches(10), Inches(0.4),
        [[(kicker.upper(), 13, BRAND_DK, True, F_HEAD, 3)]])
    txt(s, Inches(0.72), Inches(1.15), Inches(11.5), Inches(1.0),
        [[(title, 40, INK, True, F_HEAD)]])
    rect(s, Inches(0.78), Inches(2.12), Inches(1.3), Inches(0.05), BRAND_LT)
    page_footer(s, n, section)


# ============ SLIDE 1 — TONG QUAN (Cover + Overview) ============
s = slide()
bg(s, INK)
grad(s, 0, 0, EMU_W, EMU_H, INK, BRAND_DK, angle=60)
# hoa tiet tron trang tri
o = rect(s, Inches(9.4), Inches(-1.6), Inches(5.5), Inches(5.5), BRAND, MSO_SHAPE.OVAL)
o.fill.fore_color.rgb = BRAND_DK
o2 = rect(s, Inches(10.8), Inches(3.2), Inches(4.2), Inches(4.2), INK_SOFT, MSO_SHAPE.OVAL)
ring = rect(s, Inches(8.6), Inches(1.1), Inches(2.6), Inches(2.6), INK, MSO_SHAPE.OVAL)
ring.fill.background()
ring.line.color.rgb = BRAND_LT
ring.line.width = Pt(1.5)

chip(s, Inches(0.9), Inches(0.85), Inches(3.6), "MARKETING  x  TEAM SALES", BRAND, INK, 12)

txt(s, Inches(0.88), Inches(1.9), Inches(9.5), Inches(2.6),
    [[("CHƯƠNG TRÌNH", 22, BRAND_LT, True, F_HEAD, 2)],
     [("TRAINING TEAM SALES", 52, WHITE, True, F_HEAD)],
     [("Chiến dịch Seeding thương hiệu ", 22, RGBColor(0xd9,0xf2,0xf7), False, F_BODY),
      ("Kingsmen", 22, BRAND_LT, True, F_HEAD)]],
    space_after=8, line_spacing=1.02)

rect(s, Inches(0.92), Inches(4.55), Inches(1.5), Inches(0.05), BRAND_LT)

txt(s, Inches(0.9), Inches(4.9), Inches(11.0), Inches(1.4),
    [[("Marketing xây dựng nội dung & chính sách  —  Sales trực tiếp triển khai thực thi.",
       16, RGBColor(0xcf,0xe9,0xef), False, F_BODY)],
     [("Tổng quan  •  Mục tiêu  •  Chính sách  •  Hướng dẫn thực hiện",
       14, BRAND_LT, True, F_HEAD)]],
    space_after=10, line_spacing=1.1)

txt(s, Inches(0.9), Inches(6.85), Inches(11.5), Inches(0.4),
    [[("keokingsmen.com", 11, BRAND_LT, True, F_HEAD, 1),
      ("     |     Tài liệu nội bộ - Phòng Marketing", 11, RGBColor(0x9d,0xc5,0xcf), False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)


# ============ SLIDE 2 — MUC TIEU ============
s = slide()
bg(s, PAGE_BG)
header(s, "01 — Định hướng", "Mục tiêu", "MỤC TIÊU", 2)

goals = [
    ("Nhận diện thương hiệu", "Phủ sóng hình ảnh kẹo Kingsmen trên MXH & hội nhóm mục tiêu."),
    ("Xây dựng niềm tin", "Tạo review, đánh giá thật giúp khách hàng an tâm lựa chọn."),
    ("Tăng tương tác tự nhiên", "Đẩy comment, thảo luận để nội dung tiếp cận rộng hơn."),
    ("Chuyển đổi ra đơn hàng", "Biến độ quan tâm từ seeding thành khách hàng thực."),
]
gx, gy = Inches(0.78), Inches(2.6)
cw, ch = Inches(5.75), Inches(1.55)
gap = Inches(0.3)
for i, (t, d) in enumerate(goals):
    col = i % 2
    row = i // 2
    x = gx + col * (cw + gap)
    y = gy + row * (ch + Inches(0.32))
    card = rect(s, x, y, cw, ch, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, LINE, Pt(1))
    try:
        card.adjustments[0] = 0.06
    except Exception:
        pass
    rect(s, x, y, Inches(0.14), ch, BRAND, MSO_SHAPE.ROUNDED_RECTANGLE)
    num = rect(s, x + Inches(0.32), y + Inches(0.3), Inches(0.62), Inches(0.62), BRAND_BG, MSO_SHAPE.OVAL)
    ntf = num.text_frame
    ntf.paragraphs[0].alignment = PP_ALIGN.CENTER
    nr = ntf.paragraphs[0].add_run(); nr.text = f"0{i+1}"
    nr.font.size = Pt(18); nr.font.bold = True; nr.font.name = F_HEAD; nr.font.color.rgb = BRAND_DK
    txt(s, x + Inches(1.15), y + Inches(0.26), cw - Inches(1.35), ch - Inches(0.4),
        [[(t, 18, INK, True, F_HEAD)],
         [(d, 12.5, INK_MUTED, False, F_BODY)]],
        space_after=4, line_spacing=1.08)

# thanh KPI
kb = rect(s, Inches(0.78), Inches(6.28), Inches(11.75), Inches(0.5), INK, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(1.0), Inches(6.28), Inches(11.3), Inches(0.5),
    [[("ĐỊNH HƯỚNG KPI:   ", 13, BRAND_LT, True, F_HEAD),
      ("Số bài seeding  •  Lượt comment/tương tác  •  Số đơn phát sinh từ seeding mỗi tuần",
       13, WHITE, False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)


# ============ SLIDE 3 — CHINH SACH ============
s = slide()
bg(s, PAGE_BG)
header(s, "02 — Nguyên tắc", "Chính sách", "CHÍNH SÁCH", 3)

# panel trai: nen tuan thu
rect(s, Inches(0.78), Inches(2.5), Inches(7.4), Inches(4.15), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, LINE, Pt(1))
txt(s, Inches(1.1), Inches(2.78), Inches(6.9), Inches(0.5),
    [[("NGUYÊN TẮC BẮT BUỘC", 15, BRAND_DK, True, F_HEAD, 1)]])
dos = [
    "Nội dung đúng chuẩn thương hiệu: tông màu, thông điệp, hình ảnh Kingsmen.",
    "Bám sát kịch bản & nội dung do Marketing cung cấp.",
    "Seeding tự nhiên, đúng nhóm khách hàng mục tiêu.",
    "Cập nhật kết quả trung thực lên App Seeding.",
    "Bảo mật thông tin nội bộ, tài khoản & chiến dịch.",
]
yy = Inches(3.35)
for d in dos:
    dot = rect(s, Inches(1.12), yy + Inches(0.06), Inches(0.26), Inches(0.26), BRAND, MSO_SHAPE.OVAL)
    dtf = dot.text_frame; dtf.margin_top=0; dtf.margin_bottom=0
    dp = dtf.paragraphs[0]; dp.alignment=PP_ALIGN.CENTER
    dr = dp.add_run(); dr.text="✓"; dr.font.size=Pt(11); dr.font.bold=True; dr.font.name=F_HEAD; dr.font.color.rgb=WHITE
    txt(s, Inches(1.55), yy - Inches(0.03), Inches(6.4), Inches(0.6),
        [[(d, 13.5, INK_SOFT, False, F_BODY)]], line_spacing=1.05)
    yy = yy + Inches(0.62)

# panel phai: khong duoc
rect(s, Inches(8.42), Inches(2.5), Inches(4.13), Inches(4.15), INK, MSO_SHAPE.ROUNDED_RECTANGLE)
txt(s, Inches(8.72), Inches(2.78), Inches(3.6), Inches(0.5),
    [[("TUYỆT ĐỐI TRÁNH", 15, BRAND_LT, True, F_HEAD, 1)]])
donts = [
    "Spam, đăng hàng loạt gây phản cảm.",
    "Nói xấu, so sánh hạ thấp đối thủ.",
    "Tự ý thay đổi thông điệp thương hiệu.",
    "Cam kết sai lệch về sản phẩm.",
]
yy = Inches(3.35)
for d in donts:
    dot = rect(s, Inches(8.74), yy + Inches(0.05), Inches(0.26), Inches(0.26), RGBColor(0xef,0x6a,0x6a), MSO_SHAPE.OVAL)
    dtf = dot.text_frame; dtf.margin_top=0; dtf.margin_bottom=0
    dp = dtf.paragraphs[0]; dp.alignment=PP_ALIGN.CENTER
    dr = dp.add_run(); dr.text="✕"; dr.font.size=Pt(10); dr.font.bold=True; dr.font.name=F_HEAD; dr.font.color.rgb=WHITE
    txt(s, Inches(9.15), yy - Inches(0.03), Inches(3.2), Inches(0.7),
        [[(d, 13, RGBColor(0xdf,0xec,0xef), False, F_BODY)]], line_spacing=1.05)
    yy = yy + Inches(0.72)


# ============ SLIDE 4 — HUONG DAN ============
s = slide()
bg(s, PAGE_BG)
header(s, "03 — Quy trình thực hiện", "Hướng dẫn", "HƯỚNG DẪN", 4)

steps = [
    ("Nhận brief", "Nhận content, tài khoản & lịch seeding từ Marketing."),
    ("Đăng & bình luận", "Đăng bài / comment đúng kịch bản, đúng khung giờ."),
    ("Nuôi tương tác", "Tương tác chéo, trả lời, giữ tài khoản tự nhiên."),
    ("Ghi nhận", "Cập nhật kết quả & link lên App Seeding Kingsmen."),
    ("Báo cáo & tối ưu", "Tổng hợp số liệu, phản hồi để Marketing tối ưu."),
]
n = len(steps)
sx = Inches(0.78)
total_w = Inches(11.75)
cw = Emu(int((total_w - Inches(0.4) * (n - 1)) / n))
sy = Inches(2.85)
sh = Inches(3.05)
# duong noi
rect(s, sx + Inches(0.2), sy + Inches(0.42), total_w - Inches(0.4), Emu(int(Pt(2.2))), BRAND_LT)
for i, (t, d) in enumerate(steps):
    x = Emu(int(sx) + i * (int(cw) + int(Inches(0.4))))
    circ = rect(s, x + Emu(int((int(cw)-int(Inches(0.9)))/2)), sy, Inches(0.9), Inches(0.9), BRAND, MSO_SHAPE.OVAL)
    ctf = circ.text_frame; ctf.word_wrap=True; ctf.margin_top=0; ctf.margin_bottom=0
    cp = ctf.paragraphs[0]; cp.alignment=PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text=f"{i+1}"; cr.font.size=Pt(30); cr.font.bold=True; cr.font.name=F_HEAD; cr.font.color.rgb=WHITE
    card = rect(s, x, sy + Inches(1.2), cw, sh - Inches(1.2), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE, LINE, Pt(1))
    rect(s, x, sy + Inches(1.2), cw, Inches(0.12), BRAND)
    txt(s, x + Inches(0.14), sy + Inches(1.45), cw - Inches(0.28), sh - Inches(1.4),
        [[(f"BƯỚC {i+1}", 10, BRAND_DK, True, F_HEAD, 1.5)],
         [(t, 15, INK, True, F_HEAD)],
         [(d, 11.5, INK_MUTED, False, F_BODY)]],
        align=PP_ALIGN.CENTER, space_after=5, line_spacing=1.05)

txt(s, Inches(0.78), Inches(6.32), Inches(11.75), Inches(0.5),
    [[("Nguyên tắc vàng:  ", 13, BRAND_DK, True, F_HEAD),
      ("Đều đặn mỗi ngày  —  Tự nhiên như khách hàng thật  —  Ghi nhận đầy đủ để đo hiệu quả.",
       13, INK_SOFT, False, F_BODY)]],
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ============ SLIDE 5 — THANK YOU ============
s = slide()
bg(s, INK)
grad(s, 0, 0, EMU_W, EMU_H, INK, BRAND_DK, angle=60)
rect(s, Inches(9.8), Inches(-1.4), Inches(5.2), Inches(5.2), BRAND_DK, MSO_SHAPE.OVAL)
r2 = rect(s, Inches(11.3), Inches(3.6), Inches(3.8), Inches(3.8), INK_SOFT, MSO_SHAPE.OVAL)
ring = rect(s, Inches(-1.2), Inches(4.2), Inches(3.2), Inches(3.2), INK, MSO_SHAPE.OVAL)
ring.fill.background(); ring.line.color.rgb = BRAND; ring.line.width = Pt(1.2)

chip(s, Inches(0.95), Inches(1.7), Inches(2.5), "KINGSMEN TEAM", BRAND, INK, 12)
txt(s, Inches(0.9), Inches(2.55), Inches(11), Inches(2.2),
    [[("Cảm ơn!", 68, WHITE, True, F_HEAD)],
     [("Marketing đồng hành  —  Sales bứt phá doanh số.", 22, BRAND_LT, False, F_BODY)]],
    space_after=12, line_spacing=1.0)
rect(s, Inches(0.95), Inches(4.75), Inches(1.6), Inches(0.06), BRAND_LT)
txt(s, Inches(0.9), Inches(5.1), Inches(11), Inches(1.2),
    [[("Cùng nhau lan tỏa thương hiệu kẹo Kingsmen đến mọi khách hàng.",
       16, RGBColor(0xcf,0xe9,0xef), False, F_BODY)]],
    line_spacing=1.15)
txt(s, Inches(0.9), Inches(6.7), Inches(11), Inches(0.5),
    [[("keokingsmen.com", 13, BRAND_LT, True, F_HEAD, 1),
      ("     |     Phòng Marketing", 13, RGBColor(0x9d,0xc5,0xcf), False, F_BODY)]],
    anchor=MSO_ANCHOR.MIDDLE)

out = "Kingsmen-Training-Sales-Seeding.pptx"
prs.save(out)
print("SAVED:", out)
